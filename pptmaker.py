import json
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import numpy as np

# C:\Users\Ákos\Documents\GitHub\python-assessment-econ\Task1_PPTX_report\sample.json


def loader(path):  # for loading json into a list of directories representing the slides, need path for json file
    print(path)
    try:
        f = open(path)
        data = json.load(f)["presentation"]
    except:
        raise Exception("Sorry not able to use that path for making a ppt")
    return data


def ppt_processor(data):  # this function is calling the slide specific functions
    prs = Presentation()
    for i in data:  # switch for slide type
        if i['type'] == "title":
            slide_title(prs, i)
        elif i['type'] == "text":
            slide_text(prs, i)
        elif i['type'] == "list":
            slide_list(prs, i)
        elif i['type'] == "picture":
            slide_picture(prs, i)
        elif i['type'] == "plot":
            slide_plot(prs, i)
    return prs


"""Here are the functions for all the different type of slides defined in the README"""


def slide_title(prs, content):  # Title slide
    layout = prs.slide_layouts[0]  # choose a layout for slide
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title  # set the title of the slide
    title.text = content["title"]

    subtitle = slide.placeholders[1]

    subtitle.text = content["content"]
    return prs


def slide_text(prs, content):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    title.text = content["title"]

    # creating a textbox (set the position and size parameters)
    left = Inches(1)
    width = Inches(8)
    height = Inches(4)
    top = Inches(2)
    longtext = slide.shapes.add_textbox(left, top, width, height)

    longtext.text = content["content"]
    return prs


def slide_list(prs, content):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    title.text = content["title"]

    # Create bullet slide for list
    bullet = slide.shapes.placeholders[1]
    for i in content["content"]:
        p = bullet.text_frame.add_paragraph()
        p.text = i["text"]
        p.level = i["level"]
    return prs


def slide_picture(prs, content):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    title.text = content["title"]

    # add picture
    slide.shapes.add_picture(content["content"], Inches(0.5), Inches(1.75), width=Inches(9), height=Inches(5))
    return prs


def slide_plot(prs, content):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    title.text = content["title"]

    # create a chart with numpy
    arr = np.loadtxt("sample.csv", delimiter=";", dtype=float)
    chart_data = CategoryChartData()
    chart_data.categories = arr[:, 0]
    chart_data.add_series("series", arr[:, 1])

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

    slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)  # line plot
    return prs





